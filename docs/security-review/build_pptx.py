# -*- coding: utf-8 -*-
"""Build the StatsAI security-review presentation (Hebrew, RTL)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

DARK = RGBColor(0x17, 0x17, 0x17)
RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
YELLOW = RGBColor(0xCA, 0x8A, 0x04)
GREEN = RGBColor(0x05, 0x96, 0x69)
GRAY = RGBColor(0x57, 0x53, 0x4E)
LIGHT = RGBColor(0xFA, 0xFA, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
LOGO = "../../public/statsai-logo-full.png"


def set_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def add_text(slide, x, y, w, h, runs_list, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
             space_after=6):
    """runs_list: list of paragraphs; each paragraph = list of (text, size, bold, color)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for para in runs_list:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        set_rtl(p)
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Arial"
    return box


def add_bar(slide, color=RED, height=0.12):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_chip(slide, x, y, w, text, fill, font_color=WHITE, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_rtl(p)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = font_color
    return shp


def slide_header(slide, title, chip_text=None, chip_color=RED):
    add_bar(slide)
    add_text(slide, 0.6, 0.35, 12.1, 0.8, [[(title, 30, True, DARK)]])
    if chip_text:
        add_chip(slide, 0.6, 0.45, 1.7, chip_text, chip_color)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(12.1), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def bullets(slide, x, y, w, h, items, size=15, gap=10):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
            paras.append([("• " + head, size, True, DARK), (" — " + rest, size, False, GRAY)])
        else:
            paras.append([("• " + it, size, False, DARK)])
    add_text(slide, x, y, w, h, paras, space_after=gap)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_bar(s, RED, 0.18)
s.shapes.add_picture(LOGO, Inches(4.42), Inches(0.9), height=Inches(1.55))
add_text(s, 1.0, 2.8, 11.3, 1.0, [[("סקירת אבטחה ומוכנות לפרודקשן", 44, True, DARK)]], align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.75, 11.3, 0.6,
         [[("דוח ממצאים ותוכנית טיפול · אתר · מובייל · API · בסיס נתונים", 20, False, GRAY)]],
         align=PP_ALIGN.CENTER)
verdict = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.42), Inches(4.7), Inches(6.5), Inches(0.75))
verdict.fill.solid()
verdict.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
verdict.line.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
verdict.line.width = Pt(2)
tfv = verdict.text_frame
pv = tfv.paragraphs[0]
pv.alignment = PP_ALIGN.CENTER
set_rtl(pv)
rv = pv.add_run()
rv.text = "⚠️  לא מוכן לעלייה — מוכן עם תיקונים (1–2 שבועות)"
rv.font.size = Pt(20)
rv.font.bold = True
rv.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
add_text(s, 1.0, 6.3, 11.3, 0.8,
         [[("11.06.2026 · גרסה v0.13.9 · הסקירה בוצעה במצב קריאה בלבד — כל ממצא אומת מול הקוד", 13, False, GRAY)]],
         align=PP_ALIGN.CENTER)

# ---------- Slide 2: Scope ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "מה נבדק")
cards = [
    ("70", "נקודות קצה API", "כולל מתודה והרשאה לכל אחת"),
    ("72", "מודלי בסיס נתונים", "סכמה, אינדקסים, ייחודיות"),
    ("392", "קומיטים בהיסטוריה", "סריקת סודות מלאה"),
    ("11", "פאזות בדיקה", "ארכיטקטורה ועד לוגיקת ספורט"),
]
x = 0.7
for n, t, d in cards:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.95), Inches(1.7))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xE7, 0xE5, 0xE4)
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; set_rtl(p1)
    r1 = p1.add_run(); r1.text = n; r1.font.size = Pt(34); r1.font.bold = True; r1.font.color.rgb = RED
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; set_rtl(p2)
    r2 = p2.add_run(); r2.text = t; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = DARK
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; set_rtl(p3)
    r3 = p3.add_run(); r3.text = d; r3.font.size = Pt(11); r3.font.color.rgb = GRAY
    x += 3.1
bullets(s, 0.7, 3.7, 12.0, 3.3, [
    ("היקף", "אתר Next.js, אפליקציית iOS (Expo), כל ה-API, סכמת Prisma, ‏74 סקריפטים, תצורות ו-CI"),
    ("תחומים", "הזדהות והרשאות · אבטחת API · בסיס נתונים · פרונטאנד · מובייל · סודות · תלויות · לוגים · מוכנות תפעולית · לוגיקת הסטטיסטיקות"),
    ("מתודולוגיה", "שבעה תחומי סקירה במקביל; כל ממצא אומת ידנית מול הקוד (קובץ:שורה); ללא שינוי קוד"),
], size=16, gap=14)

# ---------- Slide 3: Snapshot ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "תמונת מצב")
kpis = [("4", "קריטיים", RED), ("12", "גבוהים", ORANGE), ("16", "בינוניים", YELLOW), ("19", "נמוכים/מידע", GREEN)]
x = 0.7
for n, t, c in kpis:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(2.95), Inches(1.55))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = c; card.line.width = Pt(2.5)
    tf = card.text_frame
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; set_rtl(p1)
    r1 = p1.add_run(); r1.text = n; r1.font.size = Pt(40); r1.font.bold = True; r1.font.color.rgb = c
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; set_rtl(p2)
    r2 = p2.add_run(); r2.text = t; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = DARK
    x += 3.1
add_text(s, 0.7, 3.5, 12.0, 3.5, [
    [("השורה התחתונה: ", 19, True, DARK), ("יסודות האבטחה חזקים מאוד — הבעיה היא שלמות נתונים ותפעול.", 19, False, DARK)],
    [("מה שחוסם עלייה אינו פריצה — אלא באג שמשחית סטטיסטיקות שחקנים בעריכת אירוע, מחיקה שקטה של נתונים ממוזגים במשיכה חוזרת, היעדר גיבויים על 26 שנות נתונים, וייבוא DB ללא רשת ביטחון.", 16, False, GRAY)],
    [("אפס: דליפת סודות · XSS · IDOR · mass-assignment · נתיב אדמין לא מוגן", 16, True, GREEN)],
], space_after=14)

# ---------- Slide 4: Strengths ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "מה נמצא חזק — לשמר", "✓", GREEN)
bullets(s, 0.7, 1.5, 12.0, 5.6, [
    ("הזדהות", "bcrypt(12), טוקנים נשמרים כ-hash בלבד, ביטול סשנים בשינוי סיסמה, refresh-rotation עם זיהוי שימוש חוזר"),
    ("הרשאות", "כל 21 נתיבי האדמין נבדקו — מוגנים בצד השרת; אפס IDOR; רשימות שדות מפורשות בכל עדכון"),
    ("סודות", "היסטוריית Git נקייה לחלוטין (392 קומיטים); CI סורק סודות ב-bundle המובייל"),
    ("XSS", "כל תוכן דינמי — חדשות טלגרם, תשובות AI, נתונים סרוקים — מרונדר כטקסט מוגן React"),
    ("מובייל", "refresh ב-Keychain בלבד, access בזיכרון בלבד, ATS מלא, אפס לוגים רגישים, אפס סודות ב-binary"),
    ("ארכיטקטורת מיזוג", "staging → preview → approve → execute → rollback — העיצוב הנכון; הפערים בביצוע, לא בתפיסה"),
], size=17, gap=14)

# ---------- Slide 5: C-1 ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "C-1 · עריכת אירוע משחיתה סטטיסטיקות קריירה", "CRITICAL", RED)
add_text(s, 0.7, 1.45, 12.0, 0.5, [[("src/app/api/events/route.ts:21-45", 13, True, GRAY)]], align=PP_ALIGN.LEFT)
bullets(s, 0.7, 2.0, 12.0, 3.2, [
    ("הבאג", "הקוד מציב goals = 1 במקום להוסיף 1 (חסר increment), וה-where חל על כל שורות השחקן — בכל העונות"),
    ("התוצאה", "הוספת שער דרך עורך המשחקים מאפסת מלך שערים היסטורי ל-1 שער בכל עונה; מחיקת אירוע מציבה 1-"),
    ("חשיפה", "כל עריכת אירוע מאז שהקוד עלה גרמה נזק שקט — נדרשת ביקורת נזק"),
], size=17, gap=14)
fix = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.5))
fix.fill.solid(); fix.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF9)
fix.line.color.rgb = GREEN; fix.line.width = Pt(2)
tff = fix.text_frame; tff.word_wrap = True
pf = tff.paragraphs[0]; pf.alignment = PP_ALIGN.RIGHT; set_rtl(pf)
rf = pf.add_run(); rf.text = "טיפול: תיקון ל-increment + צמצום לעונה ולמסגרת של המשחק (~10 שורות) · ביקורת נזק: goals < 0 + השוואה למקורות ושחזור"
rf.font.size = Pt(16); rf.font.bold = True; rf.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)

# ---------- Slide 6: C-2 ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "C-2 · משיכה חוזרת מוחקת שחקנים ממוזגים", "CRITICAL", RED)
add_text(s, 0.7, 1.45, 12.0, 0.5, [[("src/app/api/admin/fetch/route.ts:1636-1661", 13, True, GRAY)]], align=PP_ALIGN.LEFT)
bullets(s, 0.7, 2.0, 12.0, 3.2, [
    ("הבאג", "בסוף סנכרון שחקנים נמחק כל שחקן ללא apiFootballId ששמו אינו ברשימת ה-API — שחקני IFA ממוזגים לעולם אינם ברשימה"),
    ("התוצאה", "רענון סגל שגרתי משמיד בשקט שחקנים סרוקים + הסטטיסטיקות שלהם (cascade)"),
    ("עיקרון שנשבר", "\"מיזוג לא מוחק לעולם\" — נכון במנוע המיזוג, נשבר במסלול ה-fetch"),
], size=17, gap=14)
fix = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.2))
fix.fill.solid(); fix.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF9)
fix.line.color.rgb = GREEN; fix.line.width = Pt(2)
tff = fix.text_frame; tff.word_wrap = True
pf = tff.paragraphs[0]; pf.alignment = PP_ALIGN.RIGHT; set_rtl(pf)
rf = pf.add_run(); rf.text = "טיפול: מחיקה רק לשחקנים עם apiFootballId, או תיוג מקור לכל שחקן (~5 שורות)"
rf.font.size = Pt(16); rf.font.bold = True; rf.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)

# ---------- Slide 7: C-3 + C-4 ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "C-3 + C-4 · אין רשת ביטחון לנתונים", "CRITICAL", RED)
add_text(s, 0.7, 1.5, 12.0, 0.6, [[("C-3 · אין גיבויים אוטומטיים", 20, True, RED)]])
bullets(s, 0.7, 2.1, 12.0, 1.6, [
    "26 עונות (~13K משחקים, ~208K הרכבים) בעותק יחיד על דיסק יחיד — עלות שחזור: חודשי scraping",
    "טיפול: cron לילי pg_dump → רוטציה 7/30 יום → עותק מחוץ לשרת + תרגול שחזור אחד (~שעה)",
], size=16, gap=10)
add_text(s, 0.7, 3.9, 12.0, 0.6, [[("C-4 · ייבוא DB באדמין: מחיקת הכל לפני שחזור", 20, True, RED)]])
bullets(s, 0.7, 4.5, 12.0, 2.4, [
    "DROP CASCADE לכל הטבלאות ואז restore; רק \"FATAL\" נחשב כישלון; קובץ נטען כולו לזיכרון (OOM = השבתה באמצע מחיקה)",
    "קובץ פגום אחד = אובדן מוחלט של בסיס הנתונים, ללא דרך חזרה",
    "טיפול: snapshot אוטומטי לפני DROP · הגבלת גודל · שחזור לסכמה זמנית והחלפה אטומית (~חצי יום)",
], size=16, gap=10)

# ---------- Slide 8: High findings ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "ממצאים בחומרה גבוהה — תמצית", "HIGH", ORANGE)
bullets(s, 0.7, 1.45, 12.0, 5.7, [
    ("CSRF עוקף", "בדיקת Origin עם startsWith — דומיין מתחזה עובר (תיקון: שורה אחת)"),
    ("משחקים כפולים", "זיהוי כפילות שונה בין fetch ל-merge; כפילויות כבר קיימות בפרודקשן"),
    ("עריכת תוצאה", "לא מעדכנת טבלה וסטטיסטיקות — אין recompute בכלל"),
    ("מנוע המיזוג", "TOCTOU בין preview ל-execute · הרצה כפולה אפשרית · rollback מוחק נתונים מאוחרים"),
    ("ביצועים", "דפים ציבוריים כבדים ללא cache + אינדקסים חסרים — חישוב 26 עונות בכל בקשה"),
    ("עבודות ארוכות", "fetch בתוך בקשת HTTP (תקיעות ב-restart) · setup חוסם את האתר כולו (execSync)"),
    ("מפתחות AI", "נשמרים plaintext ב-DB — דולפים בכל ייצוא"),
    ("Rate-limit", "ניתן לזיוף דרך X-Forwarded-For — דורש אימות תצורת nginx (ידני)"),
], size=15.5, gap=10)

# ---------- Slide 9: Phases 0+1 ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "תוכנית טיפול · שלב 0 + שלב 1")
ph0 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(5.6))
ph0.fill.solid(); ph0.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xF7)
ph0.line.color.rgb = RED; ph0.line.width = Pt(2)
add_text(s, 7.0, 1.6, 5.5, 0.5, [[("שלב 0 · עצירת הנזק — יום אחד", 18, True, RED)]])
bullets(s, 7.0, 2.2, 5.5, 4.7, [
    "תיקון באג הסטטיסטיקות (C-1) + ביקורת נזק",
    "צמצום מחיקת שחקנים ב-fetch‏ (C-2)",
    "גיבוי לילי אוטומטי + תרגול שחזור (C-3)",
    "תיקון השוואת Origin‏ (CSRF)",
    "אימות X-Forwarded-For ב-nginx",
    "סגירת הרשמה חברתית כש-REGISTRATION_DISABLED",
], size=14, gap=8)
ph1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(5.9), Inches(5.6))
ph1.fill.solid(); ph1.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xF5)
ph1.line.color.rgb = ORANGE; ph1.line.width = Pt(2)
add_text(s, 0.8, 1.6, 5.5, 0.5, [[("שלב 1 · חוסמי עלייה — שבוע", 18, True, ORANGE)]])
bullets(s, 0.8, 2.2, 5.5, 4.7, [
    "חיזוק db-transfer: snapshot לפני DROP‏ (C-4)",
    "6 אינדקסים חסרים (15 דקות)",
    "caching לדפים הציבוריים הכבדים",
    "דה-דופליקציה של משחקים + recompute לטבלה",
    "הצפנת מפתחות AI ב-DB",
    "מקבץ Quick-wins (~2 שעות): תלויות מתות, SVG, באג פילטר, לוגים",
], size=14, gap=8)

# ---------- Slide 10: Phases 2+3 ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "תוכנית טיפול · שלב 2 + שלב 3")
ph2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(5.6))
ph2.fill.solid(); ph2.fill.fore_color.rgb = RGBColor(0xFE, 0xFC, 0xE8)
ph2.line.color.rgb = YELLOW; ph2.line.width = Pt(2)
add_text(s, 7.0, 1.6, 5.5, 0.5, [[("שלב 2 · עמידות — שבוע-שבועיים", 18, True, YELLOW)]])
bullets(s, 7.0, 2.2, 5.5, 4.7, [
    "מנוע מיזוג: אימות מחודש, נעילה, rollback בטוח",
    "spawn אסינכרוני במקום execSync",
    "מפתחות ייחודיות: Standing + competitionId",
    "הגבלות pagination + אינדקסי חיפוש trigram",
    "Sentry + הרחבת audit log",
    "מובייל: ניקוי cache ב-logout + סינון persister",
], size=14, gap=8)
ph3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(5.9), Inches(5.6))
ph3.fill.solid(); ph3.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF9)
ph3.line.color.rgb = GREEN; ph3.line.width = Pt(2)
add_text(s, 0.8, 1.6, 5.5, 0.5, [[("שלב 3 · אחרי העלייה", 18, True, GREEN)]])
bullets(s, 0.8, 2.2, 5.5, 4.7, [
    "איחוד לוגיקת התאמת שמות (3 עותקים → 1)",
    "פיצול route ה-fetch‏ (3,277 שורות) + job runner",
    "טבלת זהות מועדונים לקבוצות ששינו שם",
    "ביטול access tokens ב-logout-all",
    "חתימת קוד ל-OTA במובייל",
    "next/image · bcrypt/argon2",
], size=14, gap=8)

# ---------- Slide 11: Checklist ----------
s = prs.slides.add_slide(BLANK)
slide_header(s, "צ'קליסט עלייה לאוויר")
bullets(s, 0.7, 1.5, 12.0, 5.6, [
    "כל פריטי שלב 0 + שלב 1 סגורים",
    "ביקורת הנזק של C-1 הושלמה והערכים שוחזרו",
    "תרגול שחזור מגיבוי בוצע בהצלחה — גיבוי שלא שוחזר אינו גיבוי",
    "תצורת nginx אומתה (X-Forwarded-For נדרס בשרת)",
    "REGISTRATION_DISABLED מאומת בפרודקשן, כולל מסלול חברתי",
    "רוטציית מפתחות API חיצוניים (נקיים בגיט אך ותיקים בסביבה)",
    "רוטציית לוגים + התראת שטח דיסק",
    "בדיקת עומס בסיסית על דף הבית והטבלאות אחרי ה-caching",
], size=17, gap=13)

# ---------- Slide 12: Go/No-Go ----------
s = prs.slides.add_slide(BLANK)
add_bar(s, RED, 0.18)
add_text(s, 1.0, 1.3, 11.3, 1.0, [[("Go / No-Go", 40, True, DARK)]], align=PP_ALIGN.CENTER)
ng = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(2.4), Inches(7.9), Inches(1.0))
ng.fill.solid(); ng.fill.fore_color.rgb = RGBColor(0xFE, 0xE2, 0xE2)
ng.line.color.rgb = RED; ng.line.width = Pt(2.5)
tfn = ng.text_frame
pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.CENTER; set_rtl(pn)
rn = pn.add_run(); rn.text = "NO-GO כרגע — אך ורק בגלל C-1 עד C-4"
rn.font.size = Pt(24); rn.font.bold = True; rn.font.color.rgb = RGBColor(0x99, 0x1B, 0x1B)
add_text(s, 1.5, 3.8, 10.3, 1.6, [
    [("מדובר בימים, לא בחודשים: ", 18, True, DARK),
     ("C-1 הוא ~10 שורות, C-2 ~5 שורות, גיבויים — cron אחד.", 18, False, DARK)],
    [("הארכיטקטורה — staging→rollback, סטטיסטיקות בצד שרת, auth נקי — היא היסוד הנכון ואינה דורשת בנייה מחדש.", 16, False, GRAY)],
], align=PP_ALIGN.CENTER, space_after=12)
go = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(5.6), Inches(7.9), Inches(1.0))
go.fill.solid(); go.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF9)
go.line.color.rgb = GREEN; go.line.width = Pt(2.5)
tfg = go.text_frame
pg = tfg.paragraphs[0]; pg.alignment = PP_ALIGN.CENTER; set_rtl(pg)
rg = pg.add_run(); rg.text = "עם סגירת שלב 0 + שלב 1: GO ✓"
rg.font.size = Pt(24); rg.font.bold = True; rg.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)

prs.save("StatsAI-Security-Review-2026-06-11.pptx")
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
